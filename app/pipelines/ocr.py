"""
OCR module for StudyAgent.
Provides fallback OCR methods when pypdf fails to extract text.
"""
from pathlib import Path
from typing import List, Dict, Tuple
import io
import base64

try:
    from pdf2image import convert_from_path
    from PIL import Image
    import pytesseract
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

from app.config import (
    TESSERACT_CMD,
    MATHPIX_APP_ID,
    MATHPIX_APP_KEY,
    ANTHROPIC_API_KEY,
    ENABLE_MATHPIX_OCR,
    ENABLE_CLAUDE_VISION_OCR
)


def detect_poor_extraction(text: str, page_count: int = 1) -> Tuple[bool, str]:
    """
    Heuristic to detect if extraction quality is poor.

    Args:
        text: Extracted text
        page_count: Number of pages in document

    Returns:
        (needs_ocr: bool, reason: str)
    """
    from app.config import OCR_QUALITY_THRESHOLD, OCR_CHAR_DENSITY_THRESHOLD

    # Check 1: Very low character count
    avg_chars_per_page = len(text) / max(page_count, 1)
    if avg_chars_per_page < OCR_QUALITY_THRESHOLD:
        return True, f"Low text density: {avg_chars_per_page:.0f} chars/page (threshold: {OCR_QUALITY_THRESHOLD})"

    # Check 2: Low alphanumeric character ratio (suggests garbled text)
    alphanumeric = sum(c.isalnum() for c in text)
    total_chars = len(text.replace(' ', '').replace('\n', ''))

    if total_chars > 0:
        char_density = alphanumeric / total_chars
        if char_density < OCR_CHAR_DENSITY_THRESHOLD:
            return True, f"Low character quality: {char_density:.2%} alphanumeric (threshold: {OCR_CHAR_DENSITY_THRESHOLD:.0%})"

    # Check 3: Completely empty
    if not text.strip():
        return True, "No text extracted"

    return False, "Extraction quality appears good"


def extract_pdf_with_tesseract(file_path: str) -> str:
    """
    Extract text from PDF using Tesseract OCR.

    Args:
        file_path: Path to PDF file

    Returns:
        Extracted text

    Raises:
        RuntimeError: If dependencies not installed
        Exception: If OCR fails
    """
    if not PDF2IMAGE_AVAILABLE:
        raise RuntimeError(
            "pdf2image and pytesseract not installed. "
            "Run: pip install pdf2image pytesseract pillow"
        )

    # Set custom Tesseract path if configured
    if TESSERACT_CMD:
        pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD

    # Convert PDF pages to images
    try:
        images = convert_from_path(file_path, dpi=300)
    except Exception as e:
        raise RuntimeError(
            f"Failed to convert PDF to images: {e}\n"
            "Make sure poppler is installed (brew install poppler / apt install poppler-utils)"
        )

    # OCR each page
    text_parts = []
    for i, image in enumerate(images):
        try:
            page_text = pytesseract.image_to_string(image, lang='eng')
            if page_text.strip():
                text_parts.append(f"--- Page {i+1} ---\n{page_text}")
        except Exception as e:
            print(f"Warning: Tesseract failed on page {i+1}: {e}")
            continue

    if not text_parts:
        raise ValueError("Tesseract OCR extracted no text from any page")

    return '\n\n'.join(text_parts)


def extract_pdf_with_mathpix(file_path: str) -> str:
    """
    Extract text from PDF using Mathpix API (STEM-optimized).
    DISABLED BY DEFAULT - requires feature flag and API keys.

    Args:
        file_path: Path to PDF file

    Returns:
        Extracted text in Markdown format
    """
    if not ENABLE_MATHPIX_OCR:
        raise RuntimeError("Mathpix OCR is disabled. Set ENABLE_MATHPIX_OCR=true in .env")

    if not MATHPIX_APP_ID or not MATHPIX_APP_KEY:
        raise RuntimeError("Mathpix credentials not configured. Set MATHPIX_APP_ID and MATHPIX_APP_KEY in .env")

    if not PDF2IMAGE_AVAILABLE:
        raise RuntimeError(
            "pdf2image not installed. "
            "Run: pip install pdf2image pillow"
        )

    import requests

    # Convert first to images
    images = convert_from_path(file_path, dpi=300)

    text_parts = []
    for i, image in enumerate(images):
        # Convert PIL Image to bytes
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format='PNG')
        img_byte_arr = img_byte_arr.getvalue()

        # Base64 encode
        img_b64 = base64.b64encode(img_byte_arr).decode('utf-8')

        # Call Mathpix API
        response = requests.post(
            "https://api.mathpix.com/v3/text",
            headers={
                "app_id": MATHPIX_APP_ID,
                "app_key": MATHPIX_APP_KEY,
                "Content-Type": "application/json"
            },
            json={
                "src": f"data:image/png;base64,{img_b64}",
                "formats": ["text"],
                "ocr": ["math", "text"]
            }
        )

        if response.status_code == 200:
            result = response.json()
            if 'text' in result:
                text_parts.append(f"--- Page {i+1} ---\n{result['text']}")
        else:
            print(f"Warning: Mathpix failed on page {i+1}: {response.status_code}")

    if not text_parts:
        raise ValueError("Mathpix OCR extracted no text from any page")

    return '\n\n'.join(text_parts)


def extract_pdf_with_claude_vision(file_path: str) -> str:
    """
    Extract text from PDF using Claude Vision API.
    DISABLED BY DEFAULT - requires feature flag.

    Args:
        file_path: Path to PDF file

    Returns:
        Extracted text
    """
    if not ENABLE_CLAUDE_VISION_OCR:
        raise RuntimeError("Claude Vision OCR is disabled. Set ENABLE_CLAUDE_VISION_OCR=true in .env")

    if not ANTHROPIC_API_KEY:
        raise RuntimeError("Anthropic API key not configured")

    if not PDF2IMAGE_AVAILABLE:
        raise RuntimeError(
            "pdf2image not installed. "
            "Run: pip install pdf2image pillow"
        )

    from anthropic import Anthropic

    client = Anthropic(api_key=ANTHROPIC_API_KEY)

    # Convert to images
    images = convert_from_path(file_path, dpi=300)

    text_parts = []
    for i, image in enumerate(images):
        # Convert to base64
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format='PNG')
        img_byte_arr = img_byte_arr.getvalue()
        img_b64 = base64.b64encode(img_byte_arr).decode('utf-8')

        # Call Claude Vision
        response = client.messages.create(
            model="claude-3-5-sonnet-20241022",
            max_tokens=4096,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": "image/png",
                                "data": img_b64
                            }
                        },
                        {
                            "type": "text",
                            "text": "Extract all text from this document page. Preserve formatting and structure. If there are equations, convert to LaTeX."
                        }
                    ]
                }
            ]
        )

        page_text = response.content[0].text
        text_parts.append(f"--- Page {i+1} ---\n{page_text}")

    if not text_parts:
        raise ValueError("Claude Vision OCR extracted no text from any page")

    return '\n\n'.join(text_parts)


def extract_pptx(file_path: str) -> str:
    """
    Extract text and OCR images from PowerPoint files.

    Strategy:
    1. Extract all text from slides
    2. Extract images from slides
    3. OCR images using Tesseract (if available)
    4. Combine text + OCR results

    Args:
        file_path: Path to .pptx file

    Returns:
        Combined text from slides and images
    """
    from pptx import Presentation

    prs = Presentation(file_path)

    slide_texts = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_content = []
        slide_content.append(f"--- Slide {slide_num} ---")

        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text)

            # Extract text from tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        slide_content.append(row_text)

        # OCR images in slide (if Tesseract available)
        if PDF2IMAGE_AVAILABLE:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    try:
                        image = Image.open(io.BytesIO(shape.image.blob))

                        # Set custom Tesseract path if configured
                        if TESSERACT_CMD:
                            pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD

                        ocr_text = pytesseract.image_to_string(image, lang='eng')
                        if ocr_text.strip():
                            slide_content.append(f"[Image OCR]: {ocr_text}")
                    except Exception as e:
                        print(f"Warning: Could not OCR image on slide {slide_num}: {e}")

        slide_texts.append('\n'.join(slide_content))

    full_text = '\n\n'.join(slide_texts)

    if not full_text.strip():
        raise ValueError("No text could be extracted from PowerPoint")

    return full_text


def get_available_ocr_methods() -> List[str]:
    """
    Return list of available OCR methods based on configuration.

    Returns:
        List of method names: ['tesseract', 'mathpix', 'claude']
    """
    methods = []

    if PDF2IMAGE_AVAILABLE:
        methods.append('tesseract')

    if ENABLE_MATHPIX_OCR and MATHPIX_APP_ID and MATHPIX_APP_KEY:
        methods.append('mathpix')

    if ENABLE_CLAUDE_VISION_OCR and ANTHROPIC_API_KEY:
        methods.append('claude')

    return methods
